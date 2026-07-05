# -*- coding: utf-8 -*-
"""
通用文件解析 pipeline：PDF / Word / Excel / PPT → Markdown（含圖片）

流程：
  ① 依副檔名路由到對應解析器
  ② 抽文字（帶頁碼/投影片編號/工作表名）+ 抽嵌入圖片（記住位置）
  ③ 圖片過濾：太小、檔案太小、重複的（logo 每頁都出現）→ 不送 VLM
  ④ 有圖 → 呼叫 vlm_client.describe_images()（接口在 vlm_client.py）
     沒圖 → 純文字直接輸出
  ⑤ 合併成 Markdown，圖片描述插回原本的頁位置

用法：
  python doc_parser.py <檔案路徑>
輸出：
  output/<檔名>.md              最終 Markdown
  output/<檔名>_images/*.png    抽出的圖片
"""
import hashlib
import sys
import time
from dataclasses import dataclass, field
from pathlib import Path

# Windows 主控台預設 cp950，遇到特殊字元的檔名 print 會炸；
# 改成 utf-8 並以替代字元容錯，確保服務不因「印進度訊息」而失敗
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

# ---------- 圖片過濾門檻（依需求調整） ----------
# 以「像素面積」為主要判斷：PDF 匯出常把圖壓很小（幾KB），
# 但有內容的圖面積仍明顯大於裝飾 icon（25x25 那種）。
MIN_AREA = 10_000      # 約 100x100 以下視為 icon/裝飾
MIN_BYTES = 1_000      # 低於 1KB 幾乎必是純色塊/線條
MAX_CELL_ROWS = 300    # Excel 每個工作表最多輸出的列數
CLEANUP_TTL_HOURS = 24  # output/ 內超過此時數的舊結果自動刪除；設 0 停用

# ---------- 頁面路由門檻（模式3 = 整頁交給 VLM） ----------
FULLPAGE_MIN_IMAGES = 3   # 一頁的有效圖片達此數 → 版面複雜，整頁看
FULLPAGE_MAX_TEXT = 80    # 有圖且文字少於此字數 → 圖是主角，整頁看
FULLPAGE_DPI = 110        # 整頁渲染解析度（110dpi ≈ A4 寬 1150px）


@dataclass
class ImageItem:
    data: bytes
    ext: str
    page: int                 # 第幾頁 / 第幾張投影片（0 = 位置不明）
    context: str = ""         # 周圍文字，給 VLM 當提示
    width: int = 0
    height: int = 0
    filename: str = ""        # 存檔後的相對路徑
    abs_path: str = ""
    description: str = ""     # VLM 填回
    skipped: str = ""         # 被過濾的原因（空 = 要送 VLM）
    is_full_page: bool = False  # True = 整頁渲染圖（模式3），VLM 用整頁判讀 prompt


@dataclass
class Page:
    no: int                   # 頁碼 / 投影片編號
    label: str                # 顯示用，如「第 3 頁」「投影片 2」「工作表: 費用」
    text: str = ""
    images: list = field(default_factory=list)   # list[ImageItem]
    mode: int = 1             # 1=純文字 2=逐圖描述 3=整頁交VLM


# ================= ① 各格式解析器 =================

def parse_pdf(path: Path) -> list:
    import fitz
    doc = fitz.open(path)
    pages = []
    for i, pg in enumerate(doc, start=1):
        page = Page(no=i, label=f"第 {i} 頁", text=pg.get_text().strip())
        seen_xref = set()
        for img in pg.get_images(full=True):
            xref = img[0]
            if xref in seen_xref:
                continue
            seen_xref.add(xref)
            info = doc.extract_image(xref)
            page.images.append(ImageItem(
                data=info["image"], ext=info["ext"], page=i,
                context=page.text[:300],
                width=info.get("width", 0), height=info.get("height", 0)))
        pages.append(page)
    doc.close()
    return pages


def parse_pptx(path: Path) -> list:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        """展開群組，回傳 (texts, pictures)。"""
        texts, pics = [], []
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                t, p = walk(sh.shapes)
                texts += t
                pics += p
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics.append(sh)
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
        return texts, pics

    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts, pics = walk(slide.shapes)
        # 講者備註也抓
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(f"（備註：{note}）")
        page = Page(no=i, label=f"投影片 {i}", text="\n".join(texts))
        for pic in pics:
            try:
                img = pic.image
            except Exception:
                continue        # 連結圖等抓不到 blob 的跳過
            # pptx 的 width/height 是 EMU，換算成約略像素（96dpi）
            page.images.append(ImageItem(
                data=img.blob, ext=img.ext, page=i,
                context=page.text[:300],
                width=(pic.width or 0) // 9525, height=(pic.height or 0) // 9525))
        pages.append(page)
    return pages


def parse_docx(path: Path) -> list:
    """docx 沒有「頁」的概念（頁數是 Word 開檔時動態排版的），
    改用 Heading 樣式切「章節」當定位單位，對 LLM 引用更有意義。"""
    import docx
    from docx.table import Table

    doc = docx.Document(path)
    try:
        blocks = list(doc.iter_inner_content())      # 段落+表格照原文順序
    except AttributeError:                            # 舊版 python-docx 退路
        blocks = list(doc.paragraphs) + list(doc.tables)

    pages, lines = [], []
    cur = Page(no=1, label="章節: (開頭)")

    def flush(next_label=None):
        nonlocal cur, lines
        cur.text = "\n\n".join(lines)
        if cur.text.strip() or cur.images:
            pages.append(cur)
        lines = []
        if next_label is not None:
            cur = Page(no=len(pages) + 1, label=next_label)

    for blk in blocks:
        if isinstance(blk, Table):
            lines.append(_table_to_md([[c.text for c in row.cells] for row in blk.rows]))
            continue
        t = blk.text.strip()
        if not t:
            continue
        style = (blk.style.name or "").lower()
        if style.startswith("heading"):
            level = int("".join(ch for ch in style if ch.isdigit()) or 2)
            if level <= 2:
                flush(f"章節: {t}")     # 大標 → 開新章節（標題進章節標記，不重複進內文）
            else:
                lines.append("#" * min(level, 6) + " " + t)
        else:
            lines.append(t)
    flush()
    if not pages:
        pages = [Page(no=1, label="內文")]

    # docx 圖片藏在 rels，無法對應到段落位置 → 集中掛一個章節，誠實標明
    img_page = Page(no=len(pages) + 1, label="文件內圖片（位置無法對應）")
    ctx = pages[0].text[:300]
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        img_page.images.append(ImageItem(
            data=rel.target_part.blob,
            ext=rel.target_part.partname.ext.lstrip("."),
            page=img_page.no, context=ctx))
    if img_page.images:
        pages.append(img_page)
    return pages


def parse_xlsx(path: Path) -> list:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)   # data_only: 拿公式算完的值
    pages = []
    for idx, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for r in ws.iter_rows(values_only=True):
            if all(v is None for v in r):
                continue
            rows.append(["" if v is None else str(v) for v in r])
            if len(rows) >= MAX_CELL_ROWS:
                rows.append([f"…（超過 {MAX_CELL_ROWS} 列，已截斷）"])
                break
        text = _table_to_md(rows) if rows else ""
        page = Page(no=idx, label=f"工作表: {ws.title}", text=text)
        for img in getattr(ws, "_images", []):
            try:
                data = img._data()
            except Exception:
                continue
            page.images.append(ImageItem(
                data=data, ext="png", page=idx, context=f"工作表 {ws.title}"))
        pages.append(page)
    return pages


def _table_to_md(rows) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [list(r) + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(str(c).replace("\n", " ") for c in r) + " |" for r in rows]
    out.insert(1, "|" + " --- |" * width)
    return "\n".join(out)


PARSERS = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx, ".xlsm": parse_xlsx,
}


# ================= ③ 圖片過濾 =================

def filter_images(pages) -> tuple:
    """回傳（要送 VLM 的圖, 被過濾的圖）。重複圖只送第一張。"""
    seen_hash = {}
    keep, skipped = [], []
    for page in pages:
        for it in page.images:
            if len(it.data) < MIN_BYTES:
                it.skipped = "檔案太小（純色塊/線條）"
            elif it.width and it.height and it.width * it.height < MIN_AREA:
                it.skipped = f"面積太小 {it.width}x{it.height}（裝飾/icon）"
            else:
                h = hashlib.md5(it.data).hexdigest()
                if h in seen_hash:
                    it.skipped = f"重複圖片（同第 {seen_hash[h]} 頁）"
                else:
                    seen_hash[h] = it.page
            (skipped if it.skipped else keep).append(it)
    return keep, skipped


# ================= ⑤ 合併輸出 =================

def build_markdown(src: Path, pages) -> str:
    """位置標記用「明文標題」而非 HTML 註解 —— LLM 引用來源時
    才會自然說出「根據第 3 頁…」，註解是可忽略的弱訊號。"""
    out = [f"# 檔案: {src.name}", ""]
    for page in pages:
        out.append(f"## 【{page.label}】")
        if page.text:
            out.append(page.text)
        for it in page.images:
            if it.skipped:
                continue    # 被過濾的不寫進正文，避免噪音
            if it.filename:     # 有存檔（persist 模式）才放圖片連結，給前端顯圖用
                out.append(f"![{it.description[:40]}]({it.filename})")
            if it.description:
                tag = "頁面完整判讀（含圖文對應）" if it.is_full_page else "圖片內容"
                out.append(f"> {tag}：{it.description}")
        out.append("")
    return "\n".join(out)


def build_chunks(src: Path, pages) -> list:
    """每個位置單位一個 chunk，自帶出處（檔名+頁碼）。
    直接可用於 RAG 切塊或讓 LLM 標注引用來源。"""
    chunks = []
    for page in pages:
        imgs = [{"path": it.filename, "description": it.description}
                for it in page.images if not it.skipped]
        if not page.text.strip() and not imgs:
            continue
        chunks.append({
            "file": src.name, "page": page.no, "label": page.label,
            "text": page.text, "images": imgs,
        })
    return chunks


# ================= 頁面路由（決定每頁的處理模式） =================

def route_pages(pages, src: Path):
    """零成本規則決定每頁模式。整頁渲染只有 PDF 做得到（PyMuPDF），
    其他格式判為複雜時降級成模式2 並提示。"""
    can_render = src.suffix.lower() == ".pdf"
    for page in pages:
        kept = [it for it in page.images if not it.skipped]
        if not kept:
            page.mode = 1
        elif len(kept) >= FULLPAGE_MIN_IMAGES or len(page.text.strip()) < FULLPAGE_MAX_TEXT:
            if can_render:
                page.mode = 3
            else:
                page.mode = 2
                print(f"  [路由] {page.label}: 版面複雜但 {src.suffix} 無法整頁渲染，降級為逐圖模式")
        else:
            page.mode = 2

    # 模式3 的頁：渲染整頁、改掛一張整頁圖，該頁的個別圖不再重複送 VLM
    mode3 = [p for p in pages if p.mode == 3]
    if not mode3:
        return
    import fitz
    doc = fitz.open(src)
    for page in mode3:
        pix = doc[page.no - 1].get_pixmap(dpi=FULLPAGE_DPI)
        for it in page.images:
            if not it.skipped:
                it.skipped = "整頁模式已涵蓋"
        page.images.append(ImageItem(
            data=pix.tobytes("png"), ext="png", page=page.no,
            context=page.text[:500], width=pix.width, height=pix.height,
            is_full_page=True))
    doc.close()


# ================= 清理 =================

def cleanup_old_outputs():
    """惰性 TTL 清理：每次執行時刪除 output/ 內過期的舊結果
    （md、圖片資料夾一起刪，不會留下斷掉的圖片連結）。"""
    out_dir = Path("output")
    if CLEANUP_TTL_HOURS <= 0 or not out_dir.exists():
        return
    import shutil
    deadline = time.time() - CLEANUP_TTL_HOURS * 3600
    for item in out_dir.iterdir():
        try:
            if item.stat().st_mtime < deadline:
                shutil.rmtree(item) if item.is_dir() else item.unlink()
                print(f"[清理] 已刪除過期輸出: {item.name}")
        except OSError:
            pass    # 檔案被占用等情況，下次再清


# ================= 主流程 =================

def process(file_path, out_dir=Path("output"), persist=True) -> tuple:
    """處理單一檔案。回傳 (markdown字串, chunks列表)。
    persist=False（Web 一次性問答用）：不存圖、不寫 md，
    圖片以記憶體直接送 VLM，處理完不留任何檔案。"""
    src = Path(file_path)
    ext = src.suffix.lower()
    if ext not in PARSERS:
        raise ValueError(f"不支援的格式: {ext}（支援 {', '.join(PARSERS)}）")

    t0 = time.time()
    pages = PARSERS[ext](src)                       # ① ② 抽文字 + 圖
    filter_images(pages)                            # ③ 標記過濾
    route_pages(pages, src)                         # ③.5 每頁決定模式
    keep = [it for p in pages for it in p.images if not it.skipped]
    skipped = [it for p in pages for it in p.images if it.skipped]

    if keep:                                        # ④ 有圖才走 VLM
        if persist:
            out_dir = Path(out_dir)
            img_dir = out_dir / f"{src.stem}_images"
            img_dir.mkdir(parents=True, exist_ok=True)
            for n, it in enumerate(keep, start=1):
                fname = (f"p{it.page}_full.png" if it.is_full_page
                         else f"p{it.page}_img{n}.{it.ext}")
                (img_dir / fname).write_bytes(it.data)
                it.filename = f"{src.stem}_images/{fname}"
        import vlm_client
        vlm_client.describe_images(keep)            # VLM 吃 it.data，不依賴檔案
    # 沒圖 → 直接純文字輸出，完全不碰 VLM

    md = build_markdown(src, pages)                 # ⑤ 合併
    if persist:
        out_dir = Path(out_dir)
        out_dir.mkdir(parents=True, exist_ok=True)
        (out_dir / f"{src.stem}.md").write_text(md, encoding="utf-8")

    modes = [sum(1 for p in pages if p.mode == m) for m in (1, 2, 3)]
    print(f"完成: output/{src.stem}.md  （{time.time()-t0:.1f}s）"
          f" | 純文字 {modes[0]} 頁 / 逐圖 {modes[1]} 頁 / 整頁VLM {modes[2]} 頁"
          f" | VLM 請求 {len(keep)} 次")
    return md, build_chunks(src, pages)


def process_many(paths, out_dir=Path("output"), persist=True) -> tuple:
    """多檔模式：合併成一份 <document> 標籤包裹的 markdown（直接塞 prompt）。
    persist=True 時另落地 combined_for_llm.md 與 chunks.json（RAG 用）。
    回傳 (combined_markdown, all_chunks)。"""
    all_md, all_chunks = [], []
    for i, p in enumerate(paths, start=1):
        md, chunks = process(p, out_dir, persist)
        all_md.append(f'<document index="{i}" source="{Path(p).name}">\n{md}\n</document>')
        all_chunks.extend(chunks)

    combined = "<documents>\n" + "\n\n".join(all_md) + "\n</documents>\n"
    if persist:
        import json
        out_dir = Path(out_dir)
        (out_dir / "combined_for_llm.md").write_text(combined, encoding="utf-8")
        (out_dir / "chunks.json").write_text(
            json.dumps(all_chunks, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"\n合併輸出: {out_dir}/combined_for_llm.md（餵 LLM 用）"
              f" | {out_dir}/chunks.json（{len(all_chunks)} 塊，RAG 用）")
    return combined, all_chunks


def _expand(args) -> list:
    """把參數展開成檔案清單（資料夾 → 內含的支援格式檔）。"""
    files = []
    for a in args:
        p = Path(a)
        if p.is_dir():
            files += [f for f in sorted(p.iterdir()) if f.suffix.lower() in PARSERS]
        else:
            files.append(p)
    return files


if __name__ == "__main__":
    if len(sys.argv) < 2:
        raise SystemExit("用法: python doc_parser.py <檔案或資料夾...>")
    cleanup_old_outputs()
    targets = _expand(sys.argv[1:])
    if not targets:
        raise SystemExit("找不到可解析的檔案")
    try:
        if len(targets) == 1:
            process(targets[0])
        else:
            process_many(targets)
    except ValueError as e:
        raise SystemExit(str(e))
